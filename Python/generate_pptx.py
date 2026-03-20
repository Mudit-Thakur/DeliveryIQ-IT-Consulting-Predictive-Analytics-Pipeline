# -*- coding: utf-8 -*-
# ============================================================
# FILE: Python/generate_pptx.py
# PURPOSE: Generate executive slide deck from LIVE run data.
#
# ROOT CAUSE OF OLD BUG:
#   The previous version computed KPIs internally with its own
#   hardcoded logic, so it always produced the same numbers
#   regardless of what generate_data.py actually generated.
#
# FIX:
#   Every single number on every slide is now pulled from two
#   sources that are written fresh by generate_data.py:
#     1. CSV/run_metadata.json  -> all pre-computed KPIs
#     2. CSV/projects.csv       -> for top-N tables and charts
#     3. CSV/risks.csv          -> for risk breakdown table
#   Nothing is hardcoded. If the data changes, the PPTX changes.
# ============================================================

import sys
import io
import os

# ── Force UTF-8 on Windows ───────────────────────────────────
sys.stdout = io.TextIOWrapper(
    sys.stdout.buffer,
    encoding="utf-8",
    errors="replace",
    line_buffering=True,
)

import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ══════════════════════════════════════════════════════════════
# PATH RESOLUTION
# Works whether script lives in Python/ or project root
# ══════════════════════════════════════════════════════════════

_here         = os.path.dirname(os.path.abspath(__file__))
_project_root = os.path.join(_here, "..")

# Resolve CSV folder
_csv_dir = os.path.join(_project_root, "CSV")
if not os.path.isdir(_csv_dir):
    _csv_dir = os.path.join(_here, "CSV")

# Resolve Reports folder
_reports_dir = os.path.join(_project_root, "Reports")
if not os.path.isdir(_reports_dir):
    _reports_dir = os.path.join(_here, "Reports")
os.makedirs(_reports_dir, exist_ok=True)

METADATA_PATH = os.path.join(_csv_dir, "run_metadata.json")
PROJECTS_PATH = os.path.join(_csv_dir, "projects.csv")
RISKS_PATH    = os.path.join(_csv_dir, "risks.csv")
OUTPUT_PATH   = os.path.join(_reports_dir, "Executive_Summary.pptx")

# ══════════════════════════════════════════════════════════════
# STEP 1: LOAD LIVE DATA
# Everything the PPTX shows comes from these two sources.
# ══════════════════════════════════════════════════════════════

def load_live_data():
    """
    Load run_metadata.json and the fresh CSVs.
    Returns (meta dict, projects DataFrame, risks DataFrame).
    Crashes with a clear message if metadata is missing,
    because that means generate_data.py was not run first.
    """
    # ── Metadata (required) ───────────────────────────────────
    if not os.path.exists(METADATA_PATH):
        print("[FAIL] CSV/run_metadata.json not found.")
        print("       Run generate_data.py first, then re-run this script.")
        sys.exit(1)

    with open(METADATA_PATH, "r", encoding="utf-8") as f:
        meta = json.load(f)

    print("Metadata loaded:")
    print("  Run timestamp   : " + str(meta.get("run_timestamp", "?")))
    print("  Seed            : " + str(meta.get("seed", "?")))
    print("  Projects        : " + str(meta.get("num_projects", "?")))
    print("  Avg delay       : " + str(meta.get("avg_delay_days", "?")) + " days")
    print("  Forecast alerts : " + str(meta.get("forecast_alerts", "?")))

    # ── Projects CSV ──────────────────────────────────────────
    if not os.path.exists(PROJECTS_PATH):
        print("[FAIL] CSV/projects.csv not found.")
        sys.exit(1)

    projects_df = pd.read_csv(PROJECTS_PATH)

    # Remove duplicate rows that generate_data.py intentionally adds
    projects_df = projects_df.drop_duplicates(subset=["ProjectID"])

    # Make sure DelayDays is numeric
    projects_df["DelayDays"]     = pd.to_numeric(projects_df["DelayDays"],     errors="coerce")
    projects_df["ForecastDelay"] = pd.to_numeric(projects_df["ForecastDelay"], errors="coerce")
    projects_df["Budget"]        = pd.to_numeric(projects_df["Budget"],        errors="coerce")

    # ── Risks CSV ─────────────────────────────────────────────
    risks_df = pd.DataFrame()
    if os.path.exists(RISKS_PATH):
        risks_df = pd.read_csv(RISKS_PATH)

    return meta, projects_df, risks_df

# ══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════

def add_title_bar(slide, title_text, subtitle_text=""):
    """Blue header bar at the top of every slide."""
    bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(10), Inches(0.75)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 114, 198)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.25), Inches(0.05),
        Inches(9.5), Inches(0.42)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size  = Pt(16)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle_text:
        tb2 = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.47),
            Inches(9.5), Inches(0.28)
        )
        p2 = tb2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.size  = Pt(9)
        r2.font.color.rgb = RGBColor(200, 225, 255)


def add_text(slide, text, x, y, w, h,
             size=11, bold=False,
             color=(51, 51, 51),
             align=PP_ALIGN.LEFT,
             wrap=True):
    """Add a text box at exact position."""
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tb.text_frame.word_wrap = wrap
    p   = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = str(text)
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = RGBColor(*color)
    return tb


def add_rect(slide, x, y, w, h, fill_rgb, radius=False):
    """Add a filled rectangle (used for KPI cards, bars, tags)."""
    shape = slide.shapes.add_shape(
        1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.fill.background()
    return shape


def kpi_card(slide, x, y, w, h, value, label, fill_rgb):
    """
    Draw one KPI card: coloured background, big value, small label.
    The value and label are overlaid as text boxes on top of the rect.
    """
    add_rect(slide, x, y, w, h, fill_rgb)

    # Big number
    add_text(slide, value,
             x, y + 0.12, w, h * 0.55,
             size=26, bold=True,
             color=(255, 255, 255),
             align=PP_ALIGN.CENTER)

    # Small label
    add_text(slide, label,
             x, y + h - 0.42, w, 0.38,
             size=9,
             color=(230, 240, 255),
             align=PP_ALIGN.CENTER)


def delay_color(days):
    """Return RGB tuple based on delay severity."""
    if not isinstance(days, (int, float)):
        return (136, 136, 136)
    if days > 20:
        return (217, 83, 79)    # red
    if days > 10:
        return (232, 119, 34)   # orange
    return (92, 184, 92)        # green


def overdue_color(pct):
    if not isinstance(pct, (int, float)):
        return (136, 136, 136)
    if pct > 30:
        return (217, 83, 79)
    if pct > 15:
        return (232, 119, 34)
    return (92, 184, 92)


def fmt(val, suffix=""):
    """Format a number nicely; return 'N/A' if not a number."""
    if val is None:
        return "N/A"
    if isinstance(val, float):
        return "{:.1f}{}".format(val, suffix)
    return "{}{}".format(val, suffix)

# ══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# Each function receives (slide, meta, projects_df, risks_df)
# and paints one complete slide from live data.
# ══════════════════════════════════════════════════════════════

def build_title_slide(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full blue background
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0, 114, 198)
    bg.line.fill.background()

    # Accent stripe
    stripe = slide.shapes.add_shape(
        1, Inches(0), Inches(3.5), Inches(10), Inches(0.06)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(255, 200, 0)
    stripe.line.fill.background()

    add_text(slide,
             "IT Consulting Predictive Analytics",
             0.8, 1.8, 8.4, 1.2,
             size=30, bold=True,
             color=(255, 255, 255),
             align=PP_ALIGN.CENTER)

    add_text(slide,
             "Project Delivery  |  Resource Utilisation  |  Risk Management",
             0.8, 2.9, 8.4, 0.5,
             size=13,
             color=(200, 228, 255),
             align=PP_ALIGN.CENTER)

    add_text(slide,
             "Run: {}     Projects: {}     Seed: {}".format(
                 meta.get("run_timestamp", "N/A"),
                 meta.get("num_projects",   "N/A"),
                 meta.get("seed",           "N/A"),
             ),
             0.8, 6.6, 8.4, 0.5,
             size=9,
             color=(170, 210, 245),
             align=PP_ALIGN.CENTER)

    print("  Slide 1 built: Title")


def build_executive_summary(prs, meta, projects_df, risks_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(
        slide,
        "Executive Summary",
        "Run: {}  |  {} projects  |  {} clients  |  {} employees".format(
            meta.get("run_timestamp", "N/A"),
            meta.get("num_projects",  "N/A"),
            meta.get("num_clients",   "N/A"),
            meta.get("num_employees", "N/A"),
        )
    )

    # ── 4 KPI cards (row 1) ───────────────────────────────────
    avg_delay     = meta.get("avg_delay_days",  "N/A")
    pct_overdue   = meta.get("pct_overdue",     "N/A")
    fc_alerts     = meta.get("forecast_alerts", "N/A")
    high_risk     = meta.get("high_risk_count", "N/A")

    cards = [
        (fmt(avg_delay, " d"),   "Avg Delay (Days)",    delay_color(avg_delay)),
        (fmt(pct_overdue, "%"),  "% Projects Overdue",  overdue_color(pct_overdue)),
        (str(fc_alerts),         "Forecast Alerts",     (232, 119, 34)),
        (str(high_risk),         "High-Risk Items",     (217, 83, 79)),
    ]

    x_positions = [0.25, 2.68, 5.11, 7.54]
    for i, (val, lbl, color) in enumerate(cards):
        kpi_card(slide, x_positions[i], 0.9, 2.2, 1.3, val, lbl, color)

    # ── 3 more KPI cards (row 2) ──────────────────────────────
    cards2 = [
        (str(meta.get("overdue_count",   "N/A")), "Overdue Projects",     (217, 83, 79)),
        (str(meta.get("open_risk_count", "N/A")), "Open Risks",           (232, 119, 34)),
        (str(meta.get("num_employees",   "N/A")), "Employees Monitored",  (0, 114, 198)),
    ]

    x2 = [0.25, 3.48, 6.71]
    for i, (val, lbl, color) in enumerate(cards2):
        kpi_card(slide, x2[i], 2.35, 3.1, 1.1, val, lbl, color)

    # ── Top 5 delayed projects table ─────────────────────────
    add_text(slide, "Top 5 Most Delayed Projects",
             0.25, 3.6, 9.5, 0.3,
             size=10, bold=True, color=(0, 114, 198))

    top5 = (projects_df
            .dropna(subset=["DelayDays"])
            .nlargest(5, "DelayDays")
            [["ProjectID", "Sector", "ProjectType", "DelayDays", "ForecastDelay"]]
            .reset_index(drop=True))

    headers = ["Project ID", "Sector", "Project Type", "Delay (Days)", "Forecast (Days)"]
    col_x   = [0.25, 1.50, 3.20, 6.40, 7.90]
    col_w   = [1.20, 1.65, 3.10, 1.44, 1.44]

    # Header row background
    add_rect(slide, 0.25, 3.95, 9.5, 0.32, (0, 114, 198))
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_text(slide, hdr, cx, 3.97, cw, 0.28,
                 size=8, bold=True, color=(255, 255, 255))

    # Data rows
    row_colors = [(240, 247, 255), (255, 255, 255)]
    for i, row in top5.iterrows():
        ry = 4.27 + i * 0.36
        add_rect(slide, 0.25, ry, 9.5, 0.34, row_colors[i % 2])

        delay_val = row["DelayDays"]
        fc_val    = row["ForecastDelay"]

        values = [
            str(int(row["ProjectID"])),
            str(row["Sector"]),
            str(row["ProjectType"]),
            str(int(delay_val)) if pd.notna(delay_val) else "N/A",
            str(int(fc_val))    if pd.notna(fc_val)    else "N/A",
        ]

        for j, (val, cx, cw) in enumerate(zip(values, col_x, col_w)):
            # Make the delay column red if severe
            txt_color = (217, 83, 79) if (j == 3 and pd.notna(delay_val)
                                          and delay_val > 20) else (51, 51, 51)
            add_text(slide, val, cx, ry + 0.03, cw, 0.30,
                     size=8, color=txt_color)

    print("  Slide 2 built: Executive Summary  "
          "(avg_delay={}, pct_overdue={}%, alerts={})".format(
              avg_delay, pct_overdue, fc_alerts))


def build_sector_slide(prs, meta, projects_df, _risks_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Sector Performance Analysis",
                  "Average delay and overdue % by sector -- live from this run")

    # Compute fresh sector stats from the actual CSV
    sector_stats = (
        projects_df
        .dropna(subset=["DelayDays"])
        .groupby("Sector")["DelayDays"]
        .agg(mean="mean", count="count")
        .rename(columns={"mean": "AvgDelay", "count": "Projects"})
        .sort_values("AvgDelay", ascending=False)
        .reset_index()
    )

    add_text(slide, "Average Delay by Sector (days)",
             0.25, 0.9, 9.5, 0.3,
             size=11, bold=True, color=(0, 114, 198))

    if len(sector_stats) == 0:
        add_text(slide, "No sector data available.", 0.25, 1.3, 9.0, 0.4)
        return

    max_delay = max(float(sector_stats["AvgDelay"].max()), 1)
    bar_area_width = 6.5    # inches available for bars
    y_pos = 1.3

    for _, row in sector_stats.iterrows():
        avg   = float(row["AvgDelay"])
        count = int(row["Projects"])
        bar_w = max(0.05, (avg / max_delay) * bar_area_width)
        color = delay_color(avg)

        # Sector label
        add_text(slide, row["Sector"],
                 0.25, y_pos, 1.55, 0.35,
                 size=9, color=(51, 51, 51))

        # Bar
        add_rect(slide, 1.85, y_pos + 0.04, bar_w, 0.26, color)

        # Value label
        add_text(slide,
                 "{:.1f} d  ({} projects)".format(avg, count),
                 1.90 + bar_w, y_pos, 2.5, 0.35,
                 size=8, bold=True, color=(51, 51, 51))

        y_pos += 0.48

    # Worst sector callout box
    worst = sector_stats.iloc[0]
    add_rect(slide, 0.25, 6.4, 9.5, 0.7, (255, 245, 245))
    add_text(slide,
             "Most delayed sector this run: {}  --  {:.1f} day avg delay  "
             "({} projects)".format(
                 worst["Sector"],
                 float(worst["AvgDelay"]),
                 int(worst["Projects"])
             ),
             0.35, 6.48, 9.3, 0.55,
             size=10, bold=True, color=(217, 83, 79))

    print("  Slide 3 built: Sector Analysis  "
          "(worst={}, {:.1f}d)".format(
              sector_stats.iloc[0]["Sector"],
              float(sector_stats.iloc[0]["AvgDelay"])))


def build_project_type_slide(prs, meta, projects_df, _risks_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Project Type Analysis",
                  "Actual vs forecast delay by project type -- live from this run")

    type_stats = (
        projects_df
        .dropna(subset=["DelayDays", "ForecastDelay"])
        .groupby("ProjectType")
        .agg(
            AvgActual   = ("DelayDays",     "mean"),
            AvgForecast = ("ForecastDelay", "mean"),
            Count       = ("ProjectID",     "count"),
        )
        .sort_values("AvgActual", ascending=False)
        .reset_index()
    )

    add_text(slide, "Avg Delay (days):  Blue = Actual     Orange = Forecast",
             0.25, 0.9, 9.5, 0.3,
             size=9, color=(100, 100, 100))

    if len(type_stats) == 0:
        add_text(slide, "No project type data available.", 0.25, 1.3, 9.0, 0.4)
        return

    max_val = max(
        float(type_stats["AvgActual"].max()),
        float(type_stats["AvgForecast"].max()),
        1
    )
    bar_area = 5.5
    y_pos = 1.3

    for _, row in type_stats.iterrows():
        actual   = float(row["AvgActual"])
        forecast = float(row["AvgForecast"])
        count    = int(row["Count"])

        # Type label
        add_text(slide, row["ProjectType"],
                 0.25, y_pos, 2.0, 0.28,
                 size=8, color=(51, 51, 51))

        # Actual bar (blue)
        aw = max(0.05, (actual / max_val) * bar_area)
        add_rect(slide, 2.3, y_pos + 0.01, aw, 0.20, (0, 114, 198))
        add_text(slide, "{:.1f}".format(actual),
                 2.35 + aw, y_pos, 0.8, 0.22,
                 size=7, color=(0, 114, 198))

        # Forecast bar (orange)
        fw = max(0.05, (forecast / max_val) * bar_area)
        add_rect(slide, 2.3, y_pos + 0.25, fw, 0.20, (232, 119, 34))
        add_text(slide, "{:.1f}".format(forecast),
                 2.35 + fw, y_pos + 0.25, 0.8, 0.22,
                 size=7, color=(232, 119, 34))

        # Count
        add_text(slide, "n={}".format(count),
                 9.0, y_pos, 0.7, 0.48,
                 size=7, color=(150, 150, 150))

        y_pos += 0.65

    # Worst type callout
    worst = type_stats.iloc[0]
    add_rect(slide, 0.25, 6.4, 9.5, 0.7, (255, 245, 245))
    add_text(slide,
             "Highest delay type: {}  --  actual {:.1f}d  /  forecast {:.1f}d".format(
                 worst["ProjectType"],
                 float(worst["AvgActual"]),
                 float(worst["AvgForecast"])
             ),
             0.35, 6.48, 9.3, 0.55,
             size=10, bold=True, color=(217, 83, 79))

    print("  Slide 4 built: Project Type  "
          "(worst={})".format(type_stats.iloc[0]["ProjectType"]))


def build_risk_slide(prs, meta, _projects_df, risks_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Risk Analysis",
                  "Risk distribution from this run  --  live data")

    high_count   = int(meta.get("high_risk_count",  0))
    open_count   = int(meta.get("open_risk_count",  0))
    total_risks  = int(meta.get("total_risks",       0))

    # ── 3 risk KPI cards ──────────────────────────────────────
    risk_cards = [
        (str(total_risks),  "Total Risks",       (0,   114, 198)),
        (str(high_count),   "High Impact",        (217,  83,  79)),
        (str(open_count),   "Open (Unresolved)",  (232, 119,  34)),
    ]
    for i, (val, lbl, col) in enumerate(risk_cards):
        kpi_card(slide, 0.25 + i * 3.25, 0.9, 3.0, 1.2, val, lbl, col)

    # ── Impact breakdown bar chart ────────────────────────────
    add_text(slide, "Risk Count by Impact Level",
             0.25, 2.3, 9.5, 0.3,
             size=10, bold=True, color=(0, 114, 198))

    if len(risks_df) > 0:
        impact_counts = (risks_df["RiskImpact"]
                         .value_counts()
                         .reindex(["High", "Medium", "Low"], fill_value=0))

        impact_colors = {
            "High":   (217, 83,  79),
            "Medium": (232, 119, 34),
            "Low":    (92,  184, 92),
        }
        max_count = max(impact_counts.max(), 1)
        bar_area  = 6.0
        y_pos     = 2.75

        for impact, count in impact_counts.items():
            bw = max(0.05, (count / max_count) * bar_area)
            color = impact_colors.get(impact, (136, 136, 136))

            add_text(slide, impact, 0.25, y_pos, 1.3, 0.35, size=9)
            add_rect(slide, 1.65, y_pos + 0.04, bw, 0.26, color)
            add_text(slide, str(count),
                     1.70 + bw, y_pos, 0.8, 0.35,
                     size=9, bold=True, color=color)
            y_pos += 0.52

        # ── Status breakdown ──────────────────────────────────
        add_text(slide, "Risk Count by Status",
                 0.25, 4.6, 9.5, 0.3,
                 size=10, bold=True, color=(0, 114, 198))

        status_counts = risks_df["RiskStatus"].value_counts()
        status_colors = {
            "Open":      (217, 83,  79),
            "Mitigated": (232, 119, 34),
            "Closed":    (92,  184, 92),
        }
        max_s = max(status_counts.max(), 1)
        y_pos = 5.05

        for status, count in status_counts.items():
            bw    = max(0.05, (count / max_s) * 6.0)
            color = status_colors.get(status, (136, 136, 136))

            add_text(slide, status, 0.25, y_pos, 1.3, 0.35, size=9)
            add_rect(slide, 1.65, y_pos + 0.04, bw, 0.26, color)
            add_text(slide, str(count),
                     1.70 + bw, y_pos, 0.8, 0.35,
                     size=9, bold=True, color=color)
            y_pos += 0.48
    else:
        add_text(slide, "No risk data available.", 0.25, 2.75, 9.0, 0.4)

    print("  Slide 5 built: Risk Analysis  "
          "(high={}, open={})".format(high_count, open_count))


def build_recommendations_slide(prs, meta, projects_df, _risks_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Key Recommendations",
                  "Auto-generated from this run's data")

    # Build recommendations using the ACTUAL numbers from this run
    fc_alerts  = meta.get("forecast_alerts",  "N/A")
    open_risks = meta.get("open_risk_count",   "N/A")
    sector     = meta.get("top_delayed_sector","N/A")
    proj_type  = meta.get("top_delayed_type",  "N/A")
    pct_over   = meta.get("pct_overdue",       "N/A")
    avg_delay  = meta.get("avg_delay_days",    "N/A")

    recommendations = [
        (
            "CRITICAL",
            (217, 83, 79),
            (255, 245, 245),
            "{} projects forecast to breach SLA. Escalate to PMO "
            "immediately. Assign dedicated Project Managers.".format(fc_alerts)
        ),
        (
            "HIGH PRIORITY",
            (232, 119, 34),
            (255, 250, 240),
            "{} risks are still open. The {} sector has the highest average "
            "delay this run ({} days). Introduce weekly milestone check-ins.".format(
                open_risks, sector, fmt(avg_delay, " d"))
        ),
        (
            "ACTION",
            (0, 114, 198),
            (240, 247, 255),
            "{} projects are overdue (>10 days late). {} projects are the most "
            "delayed type. Standardise delivery milestones for this "
            "project type.".format(
                meta.get("overdue_count", "N/A"), proj_type)
        ),
        (
            "RESOURCE",
            (130, 80, 180),
            (248, 240, 255),
            "Review employee utilisation dashboard. Overloaded employees "
            "(>120% utilisation) directly correlate with project delays. "
            "Redistribute workloads before next sprint."
        ),
        (
            "MONITORING",
            (92, 184, 92),
            (245, 255, 245),
            "{}% of projects are currently overdue. Schedule a weekly "
            "stakeholder review using the Power BI dashboard. "
            "Automate refresh via Power BI Service.".format(fmt(pct_over))
        ),
    ]

    y_pos = 0.95
    for (tag, tag_rgb, bg_rgb, text) in recommendations:
        # Background strip
        add_rect(slide, 0.25, y_pos, 9.5, 0.88, bg_rgb)

        # Coloured tag
        add_rect(slide, 0.25, y_pos, 1.35, 0.88, tag_rgb)
        add_text(slide, tag,
                 0.28, y_pos + 0.22, 1.28, 0.44,
                 size=7, bold=True,
                 color=(255, 255, 255),
                 align=PP_ALIGN.CENTER)

        # Recommendation text
        add_text(slide, text,
                 1.68, y_pos + 0.10, 7.95, 0.70,
                 size=9, color=(51, 51, 51), wrap=True)

        y_pos += 1.02

    print("  Slide 6 built: Recommendations  "
          "(sector={}, type={})".format(sector, proj_type))


def build_footer_slide(prs, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 40, 80)
    bg.line.fill.background()

    add_text(slide, "Thank You",
             1, 2.2, 8, 1.0,
             size=36, bold=True,
             color=(255, 255, 255),
             align=PP_ALIGN.CENTER)

    add_text(slide,
             "Generated automatically by the IT Consulting Analytics Pipeline",
             1, 3.4, 8, 0.5,
             size=12,
             color=(170, 200, 240),
             align=PP_ALIGN.CENTER)

    add_text(slide,
             "Run: {}     Seed: {}     Projects: {}".format(
                 meta.get("run_timestamp", "N/A"),
                 meta.get("seed",           "N/A"),
                 meta.get("num_projects",   "N/A"),
             ),
             1, 4.1, 8, 0.4,
             size=9,
             color=(120, 160, 210),
             align=PP_ALIGN.CENTER)

    print("  Slide 7 built: Thank You / Footer")

# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════

def main():
    print("=" * 55)
    print("  GENERATE EXECUTIVE SLIDE DECK")
    print("=" * 55)

    # Load live data -- everything comes from here
    meta, projects_df, risks_df = load_live_data()

    print("")
    print("Building slides...")

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    build_title_slide(        prs, meta)
    build_executive_summary(  prs, meta, projects_df, risks_df)
    build_sector_slide(       prs, meta, projects_df, risks_df)
    build_project_type_slide( prs, meta, projects_df, risks_df)
    build_risk_slide(         prs, meta, projects_df, risks_df)
    build_recommendations_slide(prs, meta, projects_df, risks_df)
    build_footer_slide(       prs, meta)

    prs.save(OUTPUT_PATH)

    print("")
    print("Saved --> " + OUTPUT_PATH)
    print("")
    print("Slide summary:")
    print("  1  Title")
    print("  2  Executive Summary  (avg delay={} d, {}% overdue, {} alerts)".format(
        meta.get("avg_delay_days", "N/A"),
        meta.get("pct_overdue",    "N/A"),
        meta.get("forecast_alerts","N/A"),
    ))
    print("  3  Sector Analysis    (worst: {})".format(
        meta.get("top_delayed_sector", "N/A")))
    print("  4  Project Type       (worst: {})".format(
        meta.get("top_delayed_type",   "N/A")))
    print("  5  Risk Analysis      (high={}, open={})".format(
        meta.get("high_risk_count", "N/A"),
        meta.get("open_risk_count", "N/A"),
    ))
    print("  6  Recommendations    (auto-generated from live KPIs)")
    print("  7  Footer")


if __name__ == "__main__":
    main()